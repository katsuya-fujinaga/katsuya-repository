#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Marp Markdown → PPTX（個別説明会）。依存: python-pptx（_pptx_vendor 参照）"""
from __future__ import annotations

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    _here = Path(__file__).resolve().parent
    sys.path.insert(0, str(_here / "_pptx_vendor"))
    from pptx import Presentation
    from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
CASE_ROOT = HERE.parent
MD_PATH = CASE_ROOT / "ライティング" / "個別説明会_スライド.md"
OUT_PATH = CASE_ROOT / "スライド" / "愛され手相カウンセラー_個別説明会.pptx"


def strip_md(s: str) -> str:
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    return s.strip()


def is_table_row(line: str) -> bool:
    t = line.strip()
    return t.startswith("|") and t.endswith("|") and len(t) > 2


def is_table_sep(line: str) -> bool:
    t = line.strip().replace(" ", "")
    return bool(re.match(r"^\|?[-:|]+\|?$", t))


def parse_body_lines(raw_lines: list[str]) -> list[tuple[int, str]]:
    items: list[tuple[int, str]] = []
    i = 0
    while i < len(raw_lines):
        line = raw_lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("<!--"):
            i += 1
            continue

        if is_table_row(stripped) and not is_table_sep(stripped):
            rows = []
            while i < len(raw_lines) and is_table_row(raw_lines[i].strip()):
                row = raw_lines[i].strip()
                if is_table_sep(row):
                    i += 1
                    continue
                cells = [strip_md(c) for c in row.strip("|").split("|")]
                rows.append("　".join(cells))
                i += 1
            for r in rows:
                items.append((0, r))
            continue

        m = re.match(r"^(\s*)[-*]\s+(.*)$", line)
        if m:
            indent = len(m.group(1))
            lvl = 1 if indent >= 2 else 0
            items.append((lvl, strip_md(m.group(2))))
            i += 1
            continue

        m = re.match(r"^(\d+)\.\s+(.*)$", stripped)
        if m:
            items.append((0, strip_md(m.group(2))))
            i += 1
            continue

        if stripped.startswith(">"):
            items.append((0, strip_md(stripped.lstrip(">").strip())))
            i += 1
            continue

        if stripped.startswith("###"):
            items.append((0, strip_md(stripped.lstrip("#").strip())))
            i += 1
            continue

        if stripped.startswith("|"):
            i += 1
            continue

        items.append((0, strip_md(stripped)))
        i += 1

    return items


def split_slides(md: str) -> list[str]:
    text = md.lstrip("\ufeff")
    if text.startswith("---"):
        idx = text.find("\n---\n", 3)
        if idx != -1:
            text = text[idx + 5:]
    parts = re.split(r"\n---\s*\n", text)
    return [p.strip() for p in parts if p.strip()]


def parse_slide(chunk: str) -> tuple[str, list[tuple[int, str]]] | None:
    lines: list[str] = []
    for raw in chunk.splitlines():
        t = raw.strip()
        if not t or t.startswith("<!--"):
            continue
        lines.append(raw.rstrip())

    if not lines:
        return None

    def heading_level(line: str) -> int:
        s = line.strip()
        n = 0
        for c in s:
            if c == "#":
                n += 1
            else:
                break
        return n

    title_parts: list[str] = []
    i = 0
    h0 = heading_level(lines[0])
    if h0 == 1:
        while i < len(lines):
            t = lines[i].strip()
            if heading_level(t) != 1:
                break
            title_parts.append(strip_md(t[1:].lstrip()))
            i += 1
    elif h0 == 2:
        t0 = lines[0].strip()
        title_parts.append(strip_md(t0[2:].lstrip()))
        i = 1
    elif lines[0].strip().startswith("#"):
        t0 = lines[0].strip()
        title_parts.append(strip_md(t0.lstrip("#").strip()))
        i = 1

    title = "\n".join(title_parts) if title_parts else "（スライド）"
    body_lines = lines[i:]
    body = parse_body_lines(body_lines)
    return title, body


def build_pptx(slides: list[tuple[str, list[tuple[int, str]]]], out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    for title, body in slides:
        slide = prs.slides.add_slide(blank)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.35), Inches(12.3), Inches(1.35)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Yu Gothic"
        p.font.size = Pt(28)
        p.font.bold = True

        body_box = slide.shapes.add_textbox(
            Inches(0.55), Inches(1.85), Inches(12.2), Inches(5.4)
        )
        btf = body_box.text_frame
        btf.word_wrap = True

        if not body:
            p0 = btf.paragraphs[0]
            p0.text = " "
            p0.font.name = "Yu Gothic"
            p0.font.size = Pt(14)
            continue

        for idx, (lvl, text) in enumerate(body):
            if idx == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.text = text
            bp.level = min(lvl, 2)
            bp.font.name = "Yu Gothic"
            bp.font.size = Pt(13 if lvl else 15)
            bp.space_after = Pt(4)

    prs.save(out)


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    md = MD_PATH.read_text(encoding="utf-8")
    chunks = split_slides(md)
    parsed: list[tuple[str, list[tuple[int, str]]]] = []
    for ch in chunks:
        r = parse_slide(ch)
        if r:
            parsed.append(r)

    if not parsed:
        raise SystemExit("no slides parsed")

    build_pptx(parsed, OUT_PATH)
    print(f"Wrote {OUT_PATH} ({len(parsed)} slides)")


if __name__ == "__main__":
    main()
